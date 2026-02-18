#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path
from typing import List

from pptx import Presentation


def split_slides(markdown_text: str) -> List[str]:
    blocks: List[str] = []
    cur: List[str] = []
    for raw in markdown_text.splitlines():
        if raw.strip() == "---":
            if cur:
                blocks.append("\n".join(cur).strip())
                cur = []
            continue
        cur.append(raw)
    if cur:
        blocks.append("\n".join(cur).strip())
    return [b for b in blocks if b.strip()]


def parse_title_and_body(slide_md: str) -> tuple[str, List[str]]:
    lines = [ln.rstrip() for ln in slide_md.splitlines()]
    title = ""
    body: List[str] = []
    for ln in lines:
        s = ln.strip()
        if not s:
            body.append("")
            continue
        if not title and s.startswith("#"):
            title = s.lstrip("#").strip()
            continue
        body.append(ln)
    if not title:
        for ln in lines:
            s = ln.strip()
            if s:
                title = s
                break
    if not title:
        title = "Slide"
    return title, body


def markdown_to_pptx(in_md: Path, out_pptx: Path) -> None:
    text = in_md.read_text(encoding="utf-8")
    slides_md = split_slides(text)
    prs = Presentation()

    for idx, slide_md in enumerate(slides_md):
        title, body_lines = parse_title_and_body(slide_md)
        if idx == 0:
            layout = prs.slide_layouts[0]  # title slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            subtitle = slide.placeholders[1]
            subtitle.text = "\n".join([ln.strip() for ln in body_lines if ln.strip()])
            continue

        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        first = True
        for ln in body_lines:
            s = ln.strip()
            if not s:
                continue
            level = 0
            txt = s
            if s.startswith("- "):
                txt = s[2:].strip()
                level = 0
            elif s.startswith("* "):
                txt = s[2:].strip()
                level = 0
            elif s.startswith("  - "):
                txt = s[4:].strip()
                level = 1
            elif s.startswith("### "):
                txt = s[4:].strip()
                level = 0
            elif s.startswith("## "):
                txt = s[3:].strip()
                level = 0
            elif s.startswith("# "):
                txt = s[2:].strip()
                level = 0

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.level = level

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert simple markdown slides to PPTX.")
    parser.add_argument("--in-md", required=True)
    parser.add_argument("--out-pptx", required=True)
    args = parser.parse_args()

    in_md = Path(args.in_md)
    out_pptx = Path(args.out_pptx)
    markdown_to_pptx(in_md, out_pptx)
    print(f"wrote: {out_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
